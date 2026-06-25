import openpyxl
from pptx import Presentation
import pandas as pd

# Read PowerPoint file
print("=" * 80)
print("READING POWERPOINT FILE")
print("=" * 80)

prs = Presentation('การลงข้อมูลทำงาน_HEI.pptx')

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                print(text)
        if shape.has_table:
            table = shape.table
            print("\n[TABLE FOUND]")
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                print(" | ".join(row_data))

# Read Excel file
print("\n\n" + "=" * 80)
print("READING EXCEL FILE")
print("=" * 80)

try:
    df = pd.read_excel('รายชื่อพนักงาน 05_2026.xlsx')
    print(f"\nTotal rows: {len(df)}")
    print(f"Columns: {list(df.columns)}")
    print("\nFirst 10 rows:")
    print(df.head(10).to_string())
    print("\n\nAll data:")
    print(df.to_string())
except Exception as e:
    print(f"Error reading Excel: {e}")
